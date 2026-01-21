# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from pptx.dml.color import RGBColor
# from io import BytesIO
# from PIL import Image
# import requests
# from app.serp import get_image

# SLIDE_WIDTH = Inches(13.33)
# SLIDE_HEIGHT = Inches(7.5)

# def create_ppt(slides):
#     # 🔧 Load your custom template instead of a blank presentation
#     prs = Presentation("/templates.pptx")  # Replace with your actual template path

#     skipped_images = []

#     for idx, slide_data in enumerate(slides["slides"]):
#         slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use layout index or name if needed

#         if idx == 0:
#             # 🏷️ Title Slide - Centered Title
#             title_box = slide.shapes.add_textbox(
#                 Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(2)
#             )
#             title_tf = title_box.text_frame
#             title_tf.text = slide_data["title"]
#             p = title_tf.paragraphs[0]
#             p.font.size = Pt(60)
#             p.font.bold = True
#             p.alignment = PP_ALIGN.CENTER
#         else:
#             # 📝 LEFT: Title
#             title_box = slide.shapes.add_textbox(
#                 Inches(0.5), Inches(0.3), Inches(6.5), Inches(1)
#             )
#             title_tf = title_box.text_frame
#             title_tf.text = slide_data["title"]
#             title_tf.paragraphs[0].font.size = Pt(32)
#             title_tf.paragraphs[0].font.bold = True

#             # 📝 LEFT: Bullets
#             content_box = slide.shapes.add_textbox(
#                 Inches(0.5), Inches(1.4), Inches(6.5), Inches(5.5)
#             )
#             tf = content_box.text_frame
#             tf.word_wrap = True

#             for i, bullet in enumerate(slide_data["bullets"]):
#                 p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
#                 p.text = bullet
#                 p.font.size = Pt(20)
#                 p.level = 0

#             # 🖼️ RIGHT: Image with WEBP support
#             try:
#                 image_url = get_image(slide_data["title"])
#                 img_response = requests.get(image_url, timeout=10)
#                 img_bytes = img_response.content

#                 with Image.open(BytesIO(img_bytes)) as img:
#                     if img.format not in ["BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"]:
#                         with BytesIO() as converted:
#                             img.convert("RGB").save(converted, format="PNG")
#                             converted.seek(0)
#                             slide.shapes.add_picture(
#                                 converted,
#                                 Inches(7.2), Inches(1.2),
#                                 width=Inches(5.8),
#                                 height=Inches(5.8)
#                             )
#                     else:
#                         slide.shapes.add_picture(
#                             BytesIO(img_bytes),
#                             Inches(7.2), Inches(1.2),
#                             width=Inches(5.8),
#                             height=Inches(5.8)
#                         )
#             except Exception as e:
#                 skipped_images.append(slide_data["title"])

#     prs.save("output.pptx")

#     if skipped_images:
#         print("\n⚠️ Images were skipped for the following slides:")
#         for title in skipped_images:
#             print(f" - {title}")
#     else:
#         print("\n✅ All images added successfully.")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
from PIL import Image
import requests
from app.serp import get_image

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

def create_ppt(slides):
    # 🔧 Load your custom template instead of a blank presentation
    prs = Presentation("templates/template.pptx")  # Replace with your actual template path

    skipped_images = []

    for idx, slide_data in enumerate(slides["slides"]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use layout index or name if needed

        if idx == 0:
            # 🏷️ Title Slide - Centered Title
            title_box = slide.shapes.add_textbox(
                Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(2)
            )
            title_tf = title_box.text_frame
            title_tf.text = slide_data["title"]
            p = title_tf.paragraphs[0]
            p.font.size = Pt(60)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        else:
            # 📝 LEFT: Title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(6.5), Inches(1)
            )
            title_tf = title_box.text_frame
            title_tf.text = slide_data["title"]
            title_tf.paragraphs[0].font.size = Pt(32)
            title_tf.paragraphs[0].font.bold = True

            # 📝 LEFT: Bullets
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.4), Inches(6.5), Inches(5.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(slide_data["bullets"]):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = bullet
                p.font.size = Pt(20)
                p.level = 0
                p.space_after = Pt(10)  # Adds spacing between bullet points
                p.alignment = PP_ALIGN.LEFT  # Optional: ensures left alignment


            # 🖼️ RIGHT: Image with WEBP support
            try:
                image_url = get_image(slide_data["title"])
                img_response = requests.get(image_url, timeout=10)
                img_bytes = img_response.content

                with Image.open(BytesIO(img_bytes)) as img:
                    if img.format not in ["BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"]:
                        with BytesIO() as converted:
                            img.convert("RGB").save(converted, format="PNG")
                            converted.seek(0)
                            slide.shapes.add_picture(
                                converted,
                                Inches(7.2), Inches(1.2),
                                width=Inches(5.8),
                                height=Inches(5.8)
                            )
                    else:
                        slide.shapes.add_picture(
                            BytesIO(img_bytes),
                            Inches(7.2), Inches(1.2),
                            width=Inches(5.8),
                            height=Inches(5.8)
                        )
            except Exception as e:
                skipped_images.append(slide_data["title"])

    prs.save("output.pptx")

    if skipped_images:
        print("\n⚠️ Images were skipped for the following slides:")
        for title in skipped_images:
            print(f" - {title}")
    else:
        print("\n✅ All images added successfully.")